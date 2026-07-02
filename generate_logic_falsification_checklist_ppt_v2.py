"""基于《逻辑证伪应对系统 · 股票交易执行清单 v2》生成现代简约浅色 PPT。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = (
    Path(__file__).parent
    / "exported_docs"
    / "逻辑证伪应对系统_股票交易执行清单_v2_现代简约浅色版.pptx"
).resolve()

# 960 x 540 pt ≈ 13.33 x 7.5 in
W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(248, 251, 253)
INK = RGBColor(28, 42, 58)
MUTED = RGBColor(101, 116, 135)
FAINT = RGBColor(229, 238, 244)
CARD = RGBColor(255, 255, 255)
ACCENT = RGBColor(22, 177, 166)
ACCENT_2 = RGBColor(49, 109, 221)
WARN = RGBColor(228, 82, 92)
AMBER = RGBColor(230, 162, 45)
PAGE = RGBColor(160, 174, 192)


def pt(v):
    return Pt(v)


def inch(v):
    return Inches(v)


def fill_shape(shape, color, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is not None and hasattr(shape, "line"):
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    elif hasattr(shape, "line"):
        shape.line.fill.background()


def set_run(run, size=18, color=INK, bold=False):
    run.font.name = "Microsoft YaHei"
    run.font.size = pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, color, bold)
    return box


def add_rect(slide, x, y, w, h, fill=CARD, line=FAINT, rounded=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    fill_shape(shape, fill, line)
    return shape


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, inch(0.11))
    fill_shape(band, ACCENT, ACCENT)


def add_header(slide, title, subtitle=None):
    add_textbox(slide, title, inch(0.92), inch(0.67), inch(7.9), inch(0.7), 31, INK, True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0.92), inch(1.44), inch(0.75), inch(0.06))
    fill_shape(line, ACCENT, ACCENT)
    if subtitle:
        add_textbox(slide, subtitle, inch(0.92), inch(1.7), inch(7.9), inch(0.36), 14, MUTED)


def add_footer(slide, n):
    add_textbox(slide, f"{n:02d}", inch(12.28), inch(6.94), inch(0.42), inch(0.25), 10, PAGE, align=PP_ALIGN.RIGHT)


def add_card(slide, title, body, x, y, w, h, tag=None, tag_color=ACCENT):
    add_rect(slide, x, y, w, h, CARD, FAINT)
    if tag:
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + inch(0.25), y + inch(0.23), inch(0.4), inch(0.42))
        fill_shape(chip, RGBColor(233, 250, 247))
        add_textbox(slide, tag, x + inch(0.25), y + inch(0.28), inch(0.4), inch(0.25), 12, tag_color, True, PP_ALIGN.CENTER)
    add_textbox(slide, title, x + inch(0.25), y + inch(0.97), w - inch(0.5), inch(0.39), 20, INK, True)
    add_textbox(slide, body, x + inch(0.25), y + inch(1.56), w - inch(0.5), h - inch(1.75), 15, MUTED)


def add_bullet(slide, text, x, y, w, size=16, color=INK):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y + inch(0.07), inch(0.19), inch(0.19))
    fill_shape(box, CARD, ACCENT)
    add_textbox(slide, text, x + inch(0.39), y, w - inch(0.39), inch(0.39), size, color)


def add_table_row(slide, cols, y, col_widths, header=False):
    x = inch(0.92)
    row_h = inch(0.53 if header else 0.58)
    fill = RGBColor(240, 247, 250) if header else CARD
    for text, cw in zip(cols, col_widths):
        add_rect(slide, x, y, cw, row_h, fill, FAINT, rounded=False)
        add_textbox(slide, text, x + inch(0.14), y + inch(0.11), cw - inch(0.28), row_h - inch(0.16), 13 if header else 12, INK if header else MUTED, header)
        x += cw + inch(0.06)


def add_check_row(slide, idx, title, body, x, y, w, color=ACCENT):
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y + inch(0.03), inch(0.33), inch(0.33))
    fill_shape(dot, color)
    add_textbox(slide, str(idx), x, y + inch(0.08), inch(0.33), inch(0.2), 9, CARD, True, PP_ALIGN.CENTER)
    add_textbox(slide, title, x + inch(0.56), y, inch(1.56), inch(0.33), 17, INK, True)
    add_textbox(slide, body, x + inch(2.64), y + inch(0.01), w - inch(2.64), inch(0.33), 15, MUTED)


def blank_slide(pres):
    layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(layout)
    add_bg(slide)
    return slide


def build():
    pres = Presentation()
    pres.slide_width = W
    pres.slide_height = H
    n = 0

    # 1 封面
    n += 1
    s = blank_slide(pres)
    add_textbox(s, "逻辑证伪应对系统", inch(1.19), inch(1.64), inch(8.6), inch(0.72), 40, INK, True)
    add_textbox(s, "股票交易执行清单", inch(1.19), inch(2.39), inch(8.6), inch(0.72), 36, ACCENT, True)
    add_textbox(
        s,
        "不是每次看对，而是看错时有规则、亏得可控、看对时拿得住。",
        inch(1.25),
        inch(3.5),
        inch(8.6),
        inch(0.47),
        18,
        MUTED,
    )
    add_rect(s, inch(9.58), inch(1.64), inch(2.08), inch(3.19), RGBColor(232, 249, 246), RGBColor(204, 239, 235))
    add_textbox(s, "v2\n执行\n清单", inch(9.97), inch(2.33), inch(1.3), inch(1.67), 22, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(s, "把预测转化为预案，把情绪转化为规则", inch(1.25), inch(5.0), inch(7.78), inch(0.33), 15, ACCENT, True)
    add_footer(s, n)

    # 2 总原则
    n += 1
    s = blank_slide(pres)
    add_header(s, "一、交易前总原则", "RULES BEFORE TRADE")
    principles = [
        ("01", "预测 → 预案", "不凭感觉下单，买入前先写可检验条件。"),
        ("02", "情绪 → 规则", "盘中波动不替代确认，按预设规则执行。"),
        ("03", "逻辑证伪", "买入理由被推翻就退出，而非固定比例止损。"),
        ("04", "A 股 T+1", "以收盘确认为准，盘中假突破/假跌破不算数。"),
    ]
    for i, (tag, title, body) in enumerate(principles):
        x = inch(0.84) + (i % 2) * inch(5.61)
        y = inch(2.33) + (i // 2) * inch(2.33)
        add_card(s, title, body, x, y, inch(5.28), inch(1.92), tag)
    add_footer(s, n)

    # 3 Q1
    n += 1
    s = blank_slide(pres)
    add_header(s, "二、买入前 · Q1 我为什么买？", "禁止写「感觉看涨」，必须写可检验条件。")
    add_card(s, "周期条件", "月线趋势向上；周线回调至支撑位；日线放量确认。", inch(0.84), inch(2.47), inch(2.78), inch(1.89), "周期")
    add_card(s, "结构条件", "价格站上关键均线/平台/趋势线；支撑与阻力位置清晰。", inch(3.93), inch(2.47), inch(2.78), inch(1.89), "结构")
    add_card(s, "环境条件", "板块资金仍在流入；大盘无系统性风险。", inch(7.03), inch(2.47), inch(2.78), inch(1.89), "环境")
    add_rect(s, inch(0.92), inch(5.17), inch(11.33), inch(0.72), RGBColor(240, 251, 249), RGBColor(207, 239, 234))
    add_textbox(
        s,
        "买入理由汇总（一句话）：_______________________________________________",
        inch(1.25),
        inch(5.39),
        inch(10.67),
        inch(0.33),
        15,
        INK,
    )
    add_footer(s, n)

    # 4 Q2
    n += 1
    s = blank_slide(pres)
    add_header(s, "二、买入前 · Q2 证伪条件", "每条买入理由，都必须有对应、可观察、可执行的失效条件。")
    widths = [inch(2.36), inch(5.97), inch(1.53)]
    add_table_row(s, ["买入理由", "对应证伪条件", "确认周期"], inch(2.33), widths, True)
    rows = [
        ("月线趋势向上", "跌破 20 月线且收盘无法收回 → 长期趋势假设不成立", "月末收盘"),
        ("周线回调支撑", "周收盘跌破关键支撑 → 波段持有逻辑降级", "周五收盘"),
        ("日线放量确认", "放量突破后缩量跌回平台 → 突破真实性被证伪", "日收盘"),
        ("板块资金流入", "板块指数持续走弱，个股相对强度消失", "日/周收盘"),
    ]
    for i, row in enumerate(rows):
        add_table_row(s, row, inch(2.92) + i * inch(0.64), widths)
    add_textbox(
        s,
        "我的证伪条件汇总：1. ________  2. ________  3. ________",
        inch(0.92),
        inch(5.97),
        inch(11.33),
        inch(0.33),
        14,
        ACCENT,
        True,
    )
    add_footer(s, n)

    # 5 Q3 + Q4
    n += 1
    s = blank_slide(pres)
    add_header(s, "二、买入前 · Q3 / Q4 退出与持有", "错了在哪里退？对了如何加/持/止盈？")
    add_rect(s, inch(0.84), inch(2.33), inch(5.42), inch(3.89), RGBColor(255, 247, 247), FAINT)
    add_textbox(s, "Q3  如果错了，在哪里退出？", inch(1.08), inch(2.61), inch(4.72), inch(0.39), 18, WARN, True)
    for i, item in enumerate(
        [
            "关键止损位（结构支撑/均线/平台下沿）：________ 元",
            "确认周期（日/周/月收盘）：________",
            "触发后动作（减仓 / 清仓 / 停手）：________",
            "最大可接受亏损（金额或仓位比例）：________",
        ]
    ):
        add_textbox(s, item, inch(1.08), inch(3.22) + i * inch(0.61), inch(4.72), inch(0.39), 14, MUTED)

    add_rect(s, inch(6.86), inch(2.33), inch(5.42), inch(3.89), RGBColor(240, 251, 249), FAINT)
    add_textbox(s, "Q4  如果对了，如何处理？", inch(7.11), inch(2.61), inch(4.72), inch(0.39), 18, ACCENT, True)
    for i, item in enumerate(
        [
            "加仓条件（例：回踩支撑 + 放量再启）",
            "止盈条件（例：到达前高/乖离过大/结构转弱）",
            "持有条件（例：周线结构未破坏则继续持有）",
            "是否允许补仓摊平：□ 允许  □ 禁止（逻辑完全反向时）",
        ]
    ):
        add_textbox(s, item, inch(7.11), inch(3.22) + i * inch(0.61), inch(4.72), inch(0.39), 14, MUTED)
    add_footer(s, n)

    # 6 分层应对
    n += 1
    s = blank_slide(pres)
    add_header(s, "三、分层应对规则", "不要把所有波动都当成卖出信号。")
    levels = [
        ("轻微异常", "减仓观察", "日线走弱，但周线结构尚未破坏。", AMBER, RGBColor(255, 249, 238)),
        ("关键位失守", "执行止损", "收盘跌破预设支撑，原交易逻辑被证伪。", WARN, RGBColor(255, 247, 247)),
        ("逻辑完全反向", "清仓停手", "周期 + 结构 + 板块同时转弱，禁止补仓摊平。", ACCENT_2, RGBColor(242, 247, 255)),
    ]
    for i, (title, action, body, color, fill) in enumerate(levels):
        x = inch(0.92) + i * inch(4.06)
        add_rect(s, x, inch(2.47), inch(3.4), inch(3.06), fill, FAINT)
        add_textbox(s, title, x + inch(0.36), inch(2.97), inch(2.5), inch(0.36), 20, color, True)
        add_textbox(s, action, x + inch(0.36), inch(3.69), inch(2.5), inch(0.42), 27, INK, True)
        add_textbox(s, body, x + inch(0.36), inch(4.5), inch(2.58), inch(0.78), 15, MUTED)
    add_footer(s, n)

    # 7 收盘确认
    n += 1
    s = blank_slide(pres)
    add_header(s, "四、收盘确认规则", "减少噪音 — A 股 T+1 以收盘为准。")
    add_card(s, "D  日线", "当日收盘是否收回/守住关键位，避免被盘中假跌破带偏。", inch(1.0), inch(2.47), inch(2.53), inch(1.94), "D")
    add_card(s, "W  周线", "周五收盘是否守住支撑，决定波段逻辑是否延续。", inch(4.06), inch(2.47), inch(2.53), inch(1.94), "W")
    add_card(s, "M  月线", "月末是否维持趋势结构，判断长期生态是否改变。", inch(7.11), inch(2.47), inch(2.53), inch(1.94), "M")
    add_textbox(s, "本次交易主要参考周期：  □ 日线    □ 周线    □ 月线", inch(0), inch(5.56), W, inch(0.39), 16, ACCENT, True, PP_ALIGN.CENTER)
    add_footer(s, n)

    # 8 快速清单上
    n += 1
    s = blank_slide(pres)
    add_header(s, "五、快速执行清单（上）", "买入前 · 买入时")
    add_textbox(s, "【买入前】", inch(0.92), inch(2.25), inch(2.08), inch(0.33), 16, ACCENT, True)
    for i, item in enumerate(
        [
            "买入理由写成具体可检验条件，不是一句感觉",
            "每条买入理由都有对应的证伪条件",
            "已预设退出位和确认周期（日/周/月）",
            "已设计看对时的加仓/止盈/持有规则",
            "已明确最大可接受亏损",
        ]
    ):
        add_bullet(s, item, inch(0.92), inch(2.61) + i * inch(0.5), inch(10.83))
    add_textbox(s, "【买入时】", inch(0.92), inch(5.25), inch(2.08), inch(0.33), 16, ACCENT, True)
    for i, item in enumerate(
        [
            "股票代码/名称：________  买入价：________  仓位：________",
            "买入日期：________  计划持有周期：________",
        ]
    ):
        add_bullet(s, item, inch(0.92), inch(5.61) + i * inch(0.5), inch(10.83))
    add_footer(s, n)

    # 9 快速清单下
    n += 1
    s = blank_slide(pres)
    add_header(s, "五、快速执行清单（下）", "持仓中 · 收盘后复盘")
    add_textbox(s, "【持仓中 · 每日/每周检查】", inch(0.92), inch(2.25), inch(3.75), inch(0.33), 16, ACCENT, True)
    for i, item in enumerate(
        [
            "是否出现轻微异常？→ 减仓观察",
            "是否关键位失守（收盘确认）？→ 执行止损",
            "是否逻辑完全反向？→ 清仓停手，禁止摊平",
        ]
    ):
        add_bullet(s, item, inch(0.92), inch(2.61) + i * inch(0.5), inch(10.83))
    add_textbox(s, "【收盘后复盘】", inch(0.92), inch(4.25), inch(3.75), inch(0.33), 16, ACCENT, True)
    for i, item in enumerate(
        [
            "今日/本周收盘是否触发证伪条件？",
            "是否严格执行预案？（严格执行 / 偏差较大）",
            "是否有情绪化操作？（无 / 轻微 / 严重）",
            "交易结果总结：_______________________________________________",
        ]
    ):
        add_bullet(s, item, inch(0.92), inch(4.61) + i * inch(0.5), inch(10.83))
    add_footer(s, n)

    # 10 记录模板
    n += 1
    s = blank_slide(pres)
    add_header(s, "六、单笔交易记录模板", "打印或手填，形成可追溯的交易档案。")
    blocks = [
        ("买入理由", "周期 / 结构 / 环境 — 各写一条可检验条件"),
        ("证伪条件", "1. ________ → 动作：________   2. ________   3. ________"),
        ("交易参数", "买入价 / 仓位 / 止损位 / 确认周期（日·周·月）"),
        ("分层应对", "轻微异常·减仓 | 关键失守·止损 | 完全反向·清仓"),
        ("复盘", "严格执行：□是 □否   情绪化：□无 □轻微 □严重"),
    ]
    for i, (title, body) in enumerate(blocks):
        y = inch(2.33) + i * inch(0.81)
        add_rect(s, inch(0.96), y - inch(0.11), inch(10.78), inch(0.64), CARD, FAINT)
        add_check_row(s, i + 1, title, body, inch(1.25), y, inch(10.0))
    add_footer(s, n)

    # 11 结语
    n += 1
    s = blank_slide(pres)
    add_rect(s, inch(1.67), inch(1.94), inch(10.0), inch(3.06), CARD, FAINT)
    add_textbox(s, "一句话总结", inch(0), inch(2.47), W, inch(0.5), 28, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(
        s,
        "逻辑越具体，越容易知道什么时候错了；\n看错时有规则，比每次看对更重要。",
        inch(2.22),
        inch(3.17),
        inch(8.89),
        inch(1.25),
        24,
        INK,
        True,
        PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        "系统目标：看错时有规则 · 亏得可控 · 看对时拿得住",
        inch(0),
        inch(5.56),
        W,
        inch(0.39),
        16,
        MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, n)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
